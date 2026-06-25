"""python-pptx drawing primitives shared across every BVD deck.

These are deck-agnostic: text boxes, rectangles, charts, tables, headers and
footers. Nothing here knows anything about a specific dataset.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION

from .palette import DARK, LGRAY


def I(v):
    return Inches(v)


def P(v):
    return Pt(v)


def set_bg(sl, col):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = col


def txt(sl, text, x, y, w, h, size=12, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.name = "Calibri"
    r.font.size = P(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rect(sl, x, y, w, h, fill, radius=False, line=None):
    sh = sl.shapes.add_shape(5 if radius else 1, I(x), I(y), I(w), I(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = P(0.5)
    else:
        sh.line.fill.background()
    return sh


def chart(sl, ctype, series, x, y, w, h,
          legend=True, lpos=XL_LEGEND_POSITION.TOP,
          vals=True, yax=None, colors=None, ymax=None, ymin=None):
    cd = CategoryChartData()
    cd.categories = series[0][1]
    for nm, _, vv in series:
        cd.add_series(nm, tuple(vv))
    cf = sl.shapes.add_chart(ctype, I(x), I(y), I(w), I(h), cd)
    ch = cf.chart
    ch.has_legend = legend
    if legend:
        ch.legend.position = lpos
        ch.legend.include_in_layout = False
    if colors:
        for i, s in enumerate(ch.series):
            if i < len(colors):
                f = s.format.fill
                f.solid()
                f.fore_color.rgb = colors[i]
    if vals:
        for s in ch.series:
            s.data_labels.show_value = True
            s.data_labels.font.size = P(7.5)
    if yax:
        va = ch.value_axis
        va.has_title = True
        va.axis_title.text_frame.text = yax
        va.axis_title.text_frame.paragraphs[0].runs[0].font.size = P(10)
    if ymax is not None:
        ch.value_axis.maximum_scale = ymax
    if ymin is not None:
        ch.value_axis.minimum_scale = ymin
    return ch


def tbl(sl, rows, x, y, w, h, cw=None):
    t = sl.shapes.add_table(len(rows), len(rows[0]), I(x), I(y), I(w), I(h)).table
    if cw:
        for i, v in enumerate(cw):
            t.columns[i].width = I(v)
    for r, row in enumerate(rows):
        for c, (text, bold, color, bg, size, align) in enumerate(row):
            cell = t.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(text)
            run.font.bold = bold
            run.font.size = P(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"


def hdr(sl, title, sub):
    txt(sl, title, 0.5, 0.22, 9.0, 0.55, size=24, bold=True, color=DARK)
    txt(sl, sub, 0.5, 0.77, 9.0, 0.25, size=9.5, color=LGRAY)


def foot(sl, t_str):
    txt(sl, t_str, 0.5, 5.22, 9.0, 0.28, size=8.5, color=LGRAY, italic=True)
