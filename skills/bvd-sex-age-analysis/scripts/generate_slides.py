#!/usr/bin/env python3
"""
BVD Sex and Age Analysis – MVE Outbreak PowerPoint Generator  v4.0
============================================================
Changes in v4.0:
  - Unknown age shown as explicit chart bar/segment (not silently dropped)
  - Unknown sex shown as explicit chart bar/segment  
  - New data quality slide (slide 1) with completeness metrics
  - All footnotes computed from live data (no hardcoded strings)
  - Quality gates block generation on any reconciliation failure

Usage:  python generate_slides.py data.xlsx [--lang fr] [--out report.pptx]
Deps:   pip install pandas openpyxl python-pptx
"""
import argparse, re, sys, pathlib
from pathlib import Path
from datetime import datetime
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# >>> shared-imports (this block is inlined by tools/build_portable.py)
# In-repo, import the canonical shared library from src/bvd_common.
sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[3] / "src"))
from bvd_common.palette import (rgb, NAVY, NAVY_D, SLATE, DARK, LGRAY, WHITE,
                                TEAL, CORAL, AMBER, GREEN, ICE)
from bvd_common.pptx import I, P, set_bg, txt, rect, chart, tbl, hdr, foot
from bvd_common.stats import cfr_v, pct_v
# <<< shared-imports

# ── Skill-specific colours ─────────────────────────────────────────────────────
# Base neutrals (NAVY, SLATE, DARK, ...) come from bvd_common.palette.
# Below are the semantic colours specific to the MVE age/sex deck.
UNK_C=rgb("9CA3AF")                              # grey for Unknown age
CHILD_C=rgb("C0392B"); ADULT_C=rgb("0D6B7A")
CHILD_D=rgb("E8A89C"); ADULT_D=rgb("7BBFC8")
FEM_C=rgb("3266AD");   MAL_C=rgb("E07B39"); UNK_SEX=rgb("B0B8C4")
WARN=rgb("92400E");    PASS_C=rgb("14532D")

# Valid age range accepted: 0–120; above 120 (not in INVALID_AGES) triggers Gate 1
INVALID_AGES = {9999, 1999, 999}

# ── Translations ──────────────────────────────────────────────────────────────
T = {
 "en": {
  "dq_title":    "Data Quality Report",
  "dq_sub":      "Completeness and consistency checks — must be reviewed before distribution",
  "dq_total":    "Confirmed cases","dq_deaths":"Total deaths",
  "dq_unk_age":  "Unknown age","dq_unk_sex":"Unknown sex",
  "dq_inv_age":  "Invalid age codes found","dq_age_range":"Valid age range",
  "dq_geo_ok":   "Geographic coverage complete",
  "dq_rec_ok":   "Totals reconcile","dq_rec_warn":"Reconciliation issue",
  "dq_note":     "Unknown values are shown explicitly on all charts. Invalid age codes: ",
  "title_main":  "Ebola (MVE) Outbreak",
  "title_sub":   "Confirmed Cases & Deaths by Age Group and Gender",
  "title_geo":   "Democratic Republic of Congo · Ituri, Nord Kivu, Sud Kivu Provinces",
  "total_cases": "Total confirmed cases","total_deaths":"Total deaths",
  "female_cases":"Female cases","male_cases":"Male cases",
  "provinces":   "Provinces","health_zones":"Health Zones","health_areas":"Health Areas",
  "data_date":   "Data date","earliest_onset":"Earliest onset","source":"Source: DHIS2 Tracker",
  "s2_title":    "Confirmed Cases by Age Group and Gender",
  "s3_title":    "Deaths by Age Group and Gender",
  "s4_title":    "Case Fatality Rate (CFR) by Age Group",
  "s5_title":    "Children vs Adults",
  "s5_sub":      "Confirmed cases and deaths by age group (child <15 / adult ≥15),\ngender, province, and health zone. Unknown age shown explicitly.",
  "s6_title":    "Global: Cases and Deaths — Children vs Adults",
  "s7_title":    "Global: Key Metrics — Children vs Adults",
  "s8_title":    "By Province: Cases — Children vs Adults by Gender",
  "s9_title":    "By Province: Deaths — Children vs Adults by Gender",
  "s10_title":   "By Health Zone: Cases — Children vs Adults",
  "s10_sub":     "Confirmed cases  |  Zones with ≥10 cases shown",
  "s11_title":   "By Health Zone: Deaths and CFR — Children vs Adults",
  "s12_title":   "By Health Zone: Female vs Male — Cases and Deaths",
  "female":      "Female","male":"Male","unknown":"Unknown",
  "female_d":    "Female deaths","male_d":"Male deaths","unknown_d":"Unknown sex deaths",
  "child":       "Child (<15)","adult":"Adult (≥15)","unk_age":"Unknown age",
  "child_d":     "Child deaths","adult_d":"Adult deaths","unk_age_d":"Unknown age deaths",
  "child_f":"Child F","child_m":"Child M","adult_f":"Adult F","adult_m":"Adult M","unk_f":"Unk F","unk_m":"Unk M",
  "cases":"Cases","deaths":"Deaths","cfr":"CFR","cfr_pct":"CFR (%)","total":"Total",
  "province":"Province","zone":"Zone","child_t":"Child T","adult_t":"Adult T","unk_t":"Unk T",
  "cfr_crit":"≥ 40% (critical)","cfr_high":"25–39% (high)","cfr_mod":"< 25% (moderate)",
  "death_note": "Deaths identified using a combination of nature_alerte = 'Décès' and recorded death dates.",
  "prov_note":  "% = share of province total; (n) = count. Unk = unknown age. CFR = overall deaths ÷ confirmed cases.",
  "cfr_comparison":"CFR comparison","overall":"Overall",
  "confirmed_only":"Confirmed cases only",
  "confirmed_deaths":"Confirmed deaths | Deaths via nature_alerte='Décès' and recorded death dates",
  "female_cases_s12":"Female cases","male_cases_s12":"Male cases",
  "female_deaths_s12":"Female deaths","male_deaths_s12":"Male deaths",
  "unk_cases_s12":"Unknown sex cases","unk_deaths_s12":"Unknown sex deaths",
 },
 "fr": {
  "dq_title":    "Rapport de qualité des données",
  "dq_sub":      "Contrôles de complétude et de cohérence — à vérifier avant diffusion",
  "dq_total":    "Cas confirmés","dq_deaths":"Total des décès",
  "dq_unk_age":  "Âge inconnu","dq_unk_sex":"Sexe inconnu",
  "dq_inv_age":  "Codes d'âge invalides détectés","dq_age_range":"Plage d'âge valide",
  "dq_geo_ok":   "Couverture géographique complète",
  "dq_rec_ok":   "Totaux cohérents","dq_rec_warn":"Problème de cohérence",
  "dq_note":     "Les valeurs inconnues sont affichées explicitement sur tous les graphiques. Codes d'âge invalides : ",
  "title_main":  "Épidémie d'Ebola (MVE)",
  "title_sub":   "Cas confirmés et décès par groupe d'âge et sexe",
  "title_geo":   "République Démocratique du Congo · Ituri, Nord Kivu, Sud Kivu",
  "total_cases": "Total des cas confirmés","total_deaths":"Total des décès",
  "female_cases":"Cas féminins","male_cases":"Cas masculins",
  "provinces":   "Provinces","health_zones":"Zones de santé","health_areas":"Aires de santé",
  "data_date":   "Date des données","earliest_onset":"Début le plus précoce",
  "source":      "Source : DHIS2 Tracker",
  "s2_title":    "Cas confirmés par groupe d'âge et sexe",
  "s3_title":    "Décès par groupe d'âge et sexe",
  "s4_title":    "Taux de létalité (TL) par groupe d'âge",
  "s5_title":    "Enfants vs Adultes",
  "s5_sub":      "Cas confirmés et décès par groupe d'âge (enfant <15 / adulte ≥15),\nsexe, province et zone de santé. Âge inconnu affiché explicitement.",
  "s6_title":    "Global : Cas et décès — Enfants vs Adultes",
  "s7_title":    "Global : Indicateurs clés — Enfants vs Adultes",
  "s8_title":    "Par province : Cas — Enfants vs Adultes par sexe",
  "s9_title":    "Par province : Décès — Enfants vs Adultes par sexe",
  "s10_title":   "Par zone de santé : Cas — Enfants vs Adultes",
  "s10_sub":     "Cas confirmés  |  Zones avec ≥10 cas affichées",
  "s11_title":   "Par zone de santé : Décès et TL — Enfants vs Adultes",
  "s12_title":   "Par zone de santé : Féminin vs Masculin — Cas et décès",
  "female":      "Féminin","male":"Masculin","unknown":"Inconnu",
  "female_d":    "Décès féminins","male_d":"Décès masculins","unknown_d":"Décès sexe inconnu",
  "child":       "Enfant (<15)","adult":"Adulte (≥15)","unk_age":"Âge inconnu",
  "child_d":     "Décès enfants","adult_d":"Décès adultes","unk_age_d":"Décès âge inconnu",
  "child_f":"Enf F","child_m":"Enf M","adult_f":"Ad F","adult_m":"Ad M","unk_f":"Inc F","unk_m":"Inc M",
  "cases":"Cas","deaths":"Décès","cfr":"TL","cfr_pct":"TL (%)","total":"Total",
  "province":"Province","zone":"Zone","child_t":"Enf T","adult_t":"Ad T","unk_t":"Inc T",
  "cfr_crit":"≥ 40 % (critique)","cfr_high":"25–39 % (élevé)","cfr_mod":"< 25 % (modéré)",
  "death_note": "Décès identifiés par combinaison de nature_alerte = 'Décès' et des dates de décès enregistrées.",
  "prov_note":  "% = part du total provincial ; (n) = effectif. Inc = âge inconnu. TL = décès ÷ cas confirmés.",
  "cfr_comparison":"Comparaison des TL","overall":"Global",
  "confirmed_only":"Cas confirmés uniquement",
  "confirmed_deaths":"Décès confirmés | Via nature_alerte='Décès' et dates de décès enregistrées",
  "female_cases_s12":"Cas féminins","male_cases_s12":"Cas masculins",
  "female_deaths_s12":"Décès féminins","male_deaths_s12":"Décès masculins",
  "unk_cases_s12":"Cas sexe inconnu","unk_deaths_s12":"Décès sexe inconnu",
 }
}

# ── Skill-specific helpers ──────────────────────────────────────────────────────
# Generic primitives (I, P, set_bg, txt, rect, chart, tbl, hdr, foot) and the
# cfr_v / pct_v maths live in bvd_common. Only cfr_col is specific to this deck —
# it encodes the CFR severity thresholds used for colour-coding.
def cfr_col(v): return CORAL if v>=40 else (AMBER if v>=25 else TEAL)

# ── Quality gates ─────────────────────────────────────────────────────────────
def run_gates(conf,cd,gc,gd,pdata,D):
    errors=[]; warnings=[]
    extra_big={v for v in conf['age_ans'].dropna().unique() if v>120 and v not in INVALID_AGES}
    if extra_big:
        errors.append(f"Gate 1: Unexpected invalid age codes: {sorted(extra_big)}")
    for label,gc_,gd_,total_c,total_d in [("global",gc,gd,gc['total'],gd['total'])]:
        s=gc_['ct']+gc_['at']+gc_['unk']
        if s!=total_c: errors.append(f"Gate 2: Child+Adult+Unknown cases={s} ≠ total={total_c}")
        s=gd_['ct']+gd_['at']+gd_['unk']
        if s!=total_d: errors.append(f"Gate 3: Child+Adult+Unknown deaths={s} ≠ total={total_d}")
    ag_c=sum(D['tc']); unk_c=int(conf['ag'].isna().sum())
    if ag_c+unk_c!=gc['total']:
        errors.append(f"Gate 4: Age-group cases({ag_c})+unknown({unk_c})={ag_c+unk_c} ≠ total={gc['total']}")
    ag_d=sum(D['td']); unk_d=int(cd['ag'].isna().sum())
    if ag_d+unk_d!=gd['total']:
        errors.append(f"Gate 5: Age-group deaths({ag_d})+unknown({unk_d})={ag_d+unk_d} ≠ total={gd['total']}")
    if sum(p['cases']['total'] for p in pdata)!=gc['total']:
        errors.append(f"Gate 6: Province case totals don't sum to overall total")
    for p in pdata:
        c=p['cases']; s=c['ct']+c['at']+c['unk']
        if s!=c['total']:
            errors.append(f"Gate 7: Province {p['name']} child+adult+unk={s} ≠ {c['total']}")
    if sum(p['deaths']['total'] for p in pdata)!=gd['total']:
        errors.append(f"Gate 8: Province death totals don't sum to overall total")
    for p in pdata:
        d=p['deaths']; s=d['ct']+d['at']+d['unk']
        if s!=d['total']:
            errors.append(f"Gate 7b: Province {p['name']} child+adult+unk deaths={s} ≠ {d['total']}")
    for ag,tc,td in zip(D['labels'],D['tc'],D['td']):
        v=cfr_v(td,tc)
        if v<0 or v>120: errors.append(f"Gate 9: CFR {ag}={v}% out of range [0,120]")
    bad_sex=conf[~conf['sexe_clean'].isin(['Feminin','Masculin','Unknown'])]
    if len(bad_sex)>0:
        warnings.append(f"Gate 10: {len(bad_sex)} cases with unexpected sex values")
    # Verify case+death totals via gender sum
    g_case_sum = D['ftot']+D['mtot']+D['utot']
    if g_case_sum!=gc['total']:
        errors.append(f"Gate 11: F+M+Unknown cases={g_case_sum} ≠ total={gc['total']}")
    g_death_sum = D['fd_tot']+D['md_tot']+D['ud_tot']
    if g_death_sum!=gd['total']:
        errors.append(f"Gate 12: F+M+Unknown deaths={g_death_sum} ≠ total={gd['total']}")
    return errors,warnings

# ── Data extraction ───────────────────────────────────────────────────────────
def extract(filepath):
    df   = pd.read_excel(filepath)
    conf = df[df['classification_finale']=='Cas confirmé'].copy()

    # Age: invalid codes → NaN
    conf['age_clean'] = conf['age_ans'].apply(
        lambda x: np.nan if (pd.isna(x) or x in INVALID_AGES) else x)

    # Sex: normalise to Feminin/Masculin/Unknown
    conf['sexe_clean'] = conf['sexe'].apply(
        lambda x: x if x in ('Feminin','Masculin') else 'Unknown')

    conf['is_death'] = (
        (conf['nature_alerte']=='Décès') |
        conf['date_deces'].notna()       |
        conf['date_deces_alerte'].notna()
    )
    cd=conf[conf['is_death']].copy()

    # Age groups — include_lowest captures age=0
    bins  =[0,4,9,14,17,24,34,44,54,64,200]
    labels=['0-4','5-9','10-14','15-17','18-24','25-34','35-44','45-54','55-64','65+']
    conf['ag']=pd.cut(conf['age_clean'],bins=bins,labels=labels,right=True,include_lowest=True)
    cd['ag']  =pd.cut(cd['age_clean'],  bins=bins,labels=labels,right=True,include_lowest=True)

    conf['a2']=conf['age_clean'].apply(lambda x:'Unknown' if pd.isna(x) else ('Child' if x<15 else 'Adult'))
    cd['a2']  =cd['age_clean'].apply(  lambda x:'Unknown' if pd.isna(x) else ('Child' if x<15 else 'Adult'))

    # Metadata
    fname=Path(filepath).stem
    m=re.search(r'(\d{8})',fname)
    dstr=m.group(1) if m else datetime.today().strftime('%Y%m%d')
    ddate=datetime.strptime(dstr,'%Y%m%d').strftime('%d %B %Y')
    earliest="unknown"
    for col in ['date_debut_symptomes_notification','date_debut_symptomes']:
        try:
            p=pd.to_datetime(conf[col],errors='coerce').dropna()
            if len(p): earliest=p.min().strftime('%d %B %Y'); break
        except: pass

    # Age-group vectors (by cleaned sex)
    def sg(ct,row,col):
        try: return int(ct.loc[row,col])
        except: return 0

    ac =pd.crosstab(conf['ag'],conf['sexe_clean'])
    adc=pd.crosstab(cd['ag'],  cd['sexe_clean'])
    # Female/Male/Unknown per age group
    fc=[sg(ac,ag,'Feminin')  for ag in labels]
    mc=[sg(ac,ag,'Masculin') for ag in labels]
    uc=[sg(ac,ag,'Unknown')  for ag in labels]
    fd=[sg(adc,ag,'Feminin') for ag in labels]
    md=[sg(adc,ag,'Masculin')for ag in labels]
    ud=[sg(adc,ag,'Unknown') for ag in labels]
    # Age-unknown row (not in any age group)
    fc_unk=int(((conf['a2']=='Unknown')&(conf['sexe_clean']=='Feminin')).sum())
    mc_unk=int(((conf['a2']=='Unknown')&(conf['sexe_clean']=='Masculin')).sum())
    uc_unk=int(((conf['a2']=='Unknown')&(conf['sexe_clean']=='Unknown')).sum())
    fd_unk=int(((cd['a2']=='Unknown')  &(cd['sexe_clean']=='Feminin')).sum())
    md_unk=int(((cd['a2']=='Unknown')  &(cd['sexe_clean']=='Masculin')).sum())
    ud_unk=int(((cd['a2']=='Unknown')  &(cd['sexe_clean']=='Unknown')).sum())

    tc=[f+m+u for f,m,u in zip(fc,mc,uc)]
    td=[f+m+u for f,m,u in zip(fd,md,ud)]
    cfr_age=[cfr_v(d,c) for d,c in zip(td,tc)]

    # Child/Adult/Unknown × sex tracker
    def cag(df_sub):
        r={}
        for a2v,sx,k in [('Child','Feminin','cf'),('Child','Masculin','cm'),
                          ('Adult','Feminin','af'),('Adult','Masculin','am'),
                          ('Unknown','Feminin','uf'),('Unknown','Masculin','um')]:
            r[k]=int(((df_sub['a2']==a2v)&(df_sub['sexe_clean']==sx)).sum())
        r['ct']=r['cf']+r['cm']; r['at']=r['af']+r['am']
        r['unk']=int((df_sub['a2']=='Unknown').sum())
        r['unk_sex']=int((df_sub['sexe_clean']=='Unknown').sum())
        r['total']=len(df_sub)
        assert r['ct']+r['at']+r['unk']==r['total'], \
            f"cag mismatch: {r['ct']}+{r['at']}+{r['unk']}={r['ct']+r['at']+r['unk']} ≠ {r['total']}"
        return r

    gc=cag(conf); gd=cag(cd)

    pdata=[]
    for pv in conf['Province'].dropna().unique():
        sc=conf[conf['Province']==pv]; sd=cd[cd['Province']==pv]
        pdata.append({'name':pv,'cases':cag(sc),'deaths':cag(sd)})

    tz=conf['Zone_sante'].value_counts()
    tz=tz[tz>=10].head(9).index.tolist()
    zdata=[]
    for z in tz:
        sc=conf[conf['Zone_sante']==z]; sd=cd[cd['Zone_sante']==z]
        pv=sc['Province'].iloc[0] if len(sc) else ''
        zdata.append({'name':z,'prov':pv,'cases':cag(sc),'deaths':cag(sd)})

    # Gender totals (including Unknown sex)
    ftot=int((conf['sexe_clean']=='Feminin').sum())
    mtot=int((conf['sexe_clean']=='Masculin').sum())
    utot=int((conf['sexe_clean']=='Unknown').sum())
    fd_tot=int((cd['sexe_clean']=='Feminin').sum())
    md_tot=int((cd['sexe_clean']=='Masculin').sum())
    ud_tot=int((cd['sexe_clean']=='Unknown').sum())

    # Zone gender note — computed from data
    def zone_gender_note():
        male_maj=[(z['name'],z['cases']['cm']+z['cases']['am'],
                   z['cases']['cf']+z['cases']['af'])
                  for z in zdata
                  if (z['cases']['cm']+z['cases']['am'])>(z['cases']['cf']+z['cases']['af'])]
        if not male_maj: return "Females outnumber males in cases across all shown zones."
        return "Males outnumber females: "+" | ".join(f"{n} (M={m}, F={f})" for n,m,f in male_maj)+"."

    def zone_deaths_note():
        ranked=sorted([(z['name'],cfr_v(z['deaths']['total'],z['cases']['total']),z['cases']['total'])
                        for z in zdata if z['cases']['total']>0],key=lambda x:-x[1])
        return "Highest CFR: "+" | ".join(f"{n} ({c:.1f}%, n={tot})" for n,c,tot in ranked[:3])+"."

    # DQ metrics — extended for WHO-style DQ slide
    inv_codes={int(k):int(v) for k,v in
               conf['age_ans'][conf['age_ans'].isin(INVALID_AGES)].value_counts().items()}
    def _pdate(col):
        try: return pd.to_datetime(conf[col], errors='coerce')
        except: return pd.Series([pd.NaT]*len(conf))
    cutoff = pd.Timestamp('today')
    onset_dt   = _pdate('date_debut_symptomes_notification')
    outcome_dt = _pdate('date_deces_alerte')
    invest_dt  = _pdate('date_heure_investigation')
    sample_dt  = _pdate('date_prelevement')
    dq={
        'total':len(conf),'total_d':len(cd),
        'unk_age':int(conf['age_clean'].isna().sum()),
        'unk_age_d':int(cd['age_clean'].isna().sum()),
        'unk_age_pct':pct_v(int(conf['age_clean'].isna().sum()),len(conf)),
        'unk_age_d_pct':pct_v(int(cd['age_clean'].isna().sum()),len(cd)),
        'unk_sex':utot,'unk_sex_pct':pct_v(utot,len(conf)),
        'inv_codes':inv_codes,
        'age_min':int(conf['age_clean'].min()) if conf['age_clean'].notna().any() else 0,
        'age_max':int(conf['age_clean'].max()) if conf['age_clean'].notna().any() else 0,
        'miss_province':int(conf['Province'].isna().sum()),
        'miss_zone':int(conf['Zone_sante'].isna().sum()),
        'miss_aire':int(conf['Aire_sante'].isna().sum()) if 'Aire_sante' in conf.columns else 0,
        'reconciled':True,
        # completeness
        'comp_age':   int(conf['age_clean'].notna().sum()),
        'comp_sex':   int(conf['sexe'].isin(['Feminin','Masculin']).sum()),
        'comp_prov':  int(conf['Province'].notna().sum()),
        'comp_zone':  int(conf['Zone_sante'].notna().sum()),
        'comp_aire':  int(conf['Aire_sante'].notna().sum()) if 'Aire_sante' in conf.columns else len(conf),
        'comp_onset': int(onset_dt.notna().sum()),
        'comp_outcome':int(outcome_dt.notna().sum()),
        # validation flags
        'val_sex_ok':   bool(conf['sexe'].isin(['Feminin','Masculin']).all()),
        'val_prov_ok':  bool(conf['Province'].dropna().isin({'Ituri','Nord Kivu','Sud Kivu'}).all()),
        'val_zone_ok':  bool(conf['Zone_sante'].notna().all()),
        'val_aire_ok':  bool(conf['Aire_sante'].notna().all()) if 'Aire_sante' in conf.columns else True,
        'val_age_range_ok': bool(conf[conf['age_clean'].notna()]['age_clean'].between(0,120).all()),
        'val_dates_ok': bool(onset_dt.dropna().max() <= cutoff if onset_dt.notna().any() else True),
        'val_future_invest': int((invest_dt > cutoff).sum()),
        'val_future_sample': int((sample_dt > cutoff).sum()),
        'neg_age':    int((conf['age_ans'] < 0).sum()),
        'age_out_range': int(conf[conf['age_clean'].notna()]['age_clean'].gt(120).sum()),
        'dup_ids':    int(conf['numero_identification_alerte'].duplicated().sum()) if 'numero_identification_alerte' in conf.columns else 0,
        'future_dates': int((invest_dt > cutoff).sum() + (sample_dt > cutoff).sum()),
    }

    D={
        'dstr':dstr,'ddate':ddate,'earliest':earliest,
        'np':int(conf['Province'].nunique()),
        'nz':int(conf['Zone_sante'].nunique()),
        'na':int(conf['Aire_sante'].nunique()) if 'Aire_sante' in conf.columns else 0,
        'labels':labels,
        'fc':fc,'mc':mc,'uc':uc,'fc_unk':fc_unk,'mc_unk':mc_unk,'uc_unk':uc_unk,
        'fd':fd,'md':md,'ud':ud,'fd_unk':fd_unk,'md_unk':md_unk,'ud_unk':ud_unk,
        'tc':tc,'td':td,'cfr_age':cfr_age,
        'total':len(conf),'total_d':len(cd),
        'ftot':ftot,'mtot':mtot,'utot':utot,
        'fd_tot':fd_tot,'md_tot':md_tot,'ud_tot':ud_tot,
        'gc':gc,'gd':gd,'pdata':pdata,'zdata':zdata,'dq':dq,
        'zone_gender_note':zone_gender_note(),
        'zone_deaths_note':zone_deaths_note(),
    }
    return D,conf,cd

# ── Slide builders ────────────────────────────────────────────────────────────
def s0_dq(prs,D,t):
    """Slide 2 -- WHO-style Data Quality Assessment (3-column layout)."""
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,WHITE)
    dq=D["dq"]; n=dq["total"]
    SEC=rgb("1A3A5C"); OK_C=rgb("1A7A4A"); MISS_C=rgb("D4D8DD"); WARN_C=rgb("D97706")
    # Colours by severity:
    #   OK_C  (green)  = complete / valid
    #   WARN_C (amber) = incomplete but not necessarily wrong
    #   CORAL  (red)   = impossible / structural error (dup IDs, neg age, future dates, bad codes)
    #   LGRAY          = neutral / expected absence (unknown age/sex)

    # ── Navy header bar ───────────────────────────────────────────────────────
    rect(sl,0,0,10,0.56,NAVY_D)
    txt(sl,"Data Quality Assessment",0.35,0.07,7.5,0.4,size=18,bold=True,color=WHITE)
    txt(sl,f"Data date: {D['ddate']}  |  {n:,} confirmed cases  |  {t['source']}",
        0.35,0.41,9.3,0.17,size=8,color=rgb("9CB8D4"))

    # ── Column section headers ─────────────────────────────────────────────────
    def sec_hdr(x,y,w,label):
        rect(sl,x,y,w,0.23,rgb("E8EEF4"))
        txt(sl,label,x+0.1,y+0.03,w-0.15,0.19,size=8.5,bold=True,color=SEC)

    sec_hdr(0.28,0.66,4.45,"COMPLETENESS")
    sec_hdr(4.88,0.66,2.52,"DATA VALIDATION")
    sec_hdr(7.56,0.66,2.18,"DATA PROFILE")

    # ── COLUMN 1: Completeness bars ───────────────────────────────────────────
    # Outcome date: not a quality error in active outbreak — use neutral amber
    comp_items=[
        ("Age recorded",           dq["comp_age"],    n, False),
        ("Sex recorded",           dq["comp_sex"],    n, False),
        ("Province assigned",      dq["comp_prov"],   n, False),
        ("Health Zone assigned",   dq["comp_zone"],   n, False),
        ("Health Area assigned",   dq["comp_aire"],   n, False),
        ("Date of symptom onset",  dq["comp_onset"],  n, False),
        # outbreak context: missing = not yet resolved, not data error
        ("Final outcome recorded", dq["comp_outcome"],n, True),
    ]
    LBL_X=0.28; LBL_W=1.52; BAR_X=1.85; BAR_W=2.3; PCT_X=4.2
    ROW_H=0.31; CY=0.96

    for label,comp,total,is_surv_context in comp_items:
        pct=comp/total if total>0 else 0
        raw_pct=pct*100
        comp_pct=int(raw_pct) if raw_pct==int(raw_pct) else round(raw_pct,1)
        raw_miss=(1-pct)*100
        miss_pct=int(raw_miss) if raw_miss==int(raw_miss) else round(raw_miss,1)
        miss=total-comp

        txt(sl,label,LBL_X,CY+0.03,LBL_W,0.22,size=8,color=DARK)
        rect(sl,BAR_X,CY+0.05,BAR_W,0.17,MISS_C)

        if is_surv_context:
            # Outcome in active outbreak: amber even when low — not an error
            fill_col=OK_C if pct>=0.95 else WARN_C
        else:
            fill_col=OK_C if pct>=0.95 else (WARN_C if pct>=0.80 else CORAL)

        fill_w=BAR_W*pct
        if fill_w>0.01:
            rect(sl,BAR_X,CY+0.05,fill_w,0.17,fill_col)

        txt(sl,f"{comp_pct}%",PCT_X+0.04,CY+0.03,0.52,0.22,size=8,bold=True,
            color=fill_col if pct<1.0 else OK_C)

        CY+=ROW_H

    # ── COLUMN 2: Validation checklist ───────────────────────────────────────
    val_items=[
        (dq["val_sex_ok"],         "Valid sex codes"),
        (dq["val_prov_ok"],        "Valid province codes"),
        (dq["val_zone_ok"],        "Valid Health Zone codes"),
        (dq["val_aire_ok"],        "Valid Health Area codes"),
        (dq["val_dates_ok"],       "No future notification dates"),
        (dq["val_age_range_ok"],   "Valid age range (0-120)"),
        (dq["dup_ids"]==0,         "No duplicate case IDs"),
        (dq["neg_age"]==0,         "No negative ages"),
        (True,                     "Invalid age codes reclassified"),
    ]
    VY=0.96; ICON_X=4.88; LBL_VX=5.14
    for ok,label in val_items:
        ic=OK_C if ok else WARN_C
        bg=rgb("DCFCE7") if ok else rgb("FEF3C7")
        rect(sl,ICON_X,VY+0.02,0.2,0.2,bg,radius=True)
        symbol="+" if ok else "!"
        txt(sl,symbol,ICON_X+0.01,VY+0.01,0.19,0.22,size=9,bold=True,color=ic,
            align=PP_ALIGN.CENTER)
        txt(sl,label,LBL_VX,VY+0.03,2.25,0.21,size=8,
            color=DARK if ok else rgb("92400E"))
        VY+=0.26

    # ── COLUMN 3: Data Profile ─────────────────────────────────────────────────
    # Order: unknown values → data integrity → structural errors
    PX=7.56; PY=0.96; VW=2.18

    def prow(lbl,val,color=None):
        nonlocal PY
        vc=color if color else DARK
        txt(sl,lbl,PX,PY,1.35,0.2,size=8,color=LGRAY)
        txt(sl,str(val),PX+1.35,PY,0.78,0.2,size=8,bold=True,color=vc,
            align=PP_ALIGN.RIGHT)
        PY+=0.215

    # Unknown values — amber (incomplete, not erroneous)
    prow("Unknown age",  f"{dq['unk_age']} cases",  WARN_C if dq['unk_age']>0 else OK_C)
    prow("Unknown sex",  f"{dq['unk_sex']} cases",  WARN_C if dq['unk_sex']>0 else OK_C)
    PY+=0.04

    # Data integrity — red if non-zero
    prow("Future dates",      dq["future_dates"],   CORAL if dq["future_dates"]>0  else OK_C)
    prow("Missing geography", dq["miss_province"]+dq["miss_zone"],
                               CORAL if dq["miss_province"]+dq["miss_zone"]>0 else OK_C)
    prow("Duplicate case IDs",dq["dup_ids"],         CORAL if dq["dup_ids"]>0      else OK_C)
    PY+=0.04

    # Recoded values — amber (handled, not errors)
    txt(sl,"Recoded invalid age values",PX,PY,VW,0.2,size=8,bold=True,color=SEC); PY+=0.22
    for code,cnt in sorted(dq["inv_codes"].items()):
        txt(sl,f"  {code}",PX,PY,0.85,0.2,size=8,color=LGRAY)
        txt(sl,f"({cnt})",PX+0.85,PY,1.28,0.2,size=8,bold=True,color=WARN_C,
            align=PP_ALIGN.RIGHT)
        PY+=0.2
    txt(sl,"-> classified as Unknown age",PX+0.1,PY,VW-0.1,0.2,size=7.5,
        color=LGRAY,italic=True); PY+=0.24

    # Structural impossibles — red if non-zero
    prow("Negative ages",      dq["neg_age"],        CORAL if dq["neg_age"]>0      else OK_C)
    prow("Age outside 0-120",  dq["age_out_range"],  CORAL if dq["age_out_range"]>0 else OK_C)

    # ── Vertical dividers ─────────────────────────────────────────────────────
    for xd in [4.73, 7.42]:
        rect(sl,xd,0.64,0.015,4.7,rgb("DDE3EA"))

    # ── Footer ────────────────────────────────────────────────────────────────
    rect(sl,0,5.28,10,0.34,rgb("F0F4F8"))
    footer=(
        "Data quality checks are performed on the raw DHIS2 Tracker extract "
        "prior to generating epidemiological analyses. "
        "Invalid values are retained where appropriate and explicitly classified "
        "as Unknown rather than excluded from the analysis."
    )
    txt(sl,footer,0.3,5.31,9.4,0.28,size=7.5,color=rgb("4A5568"),italic=True)


def s1_title(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,NAVY_D)
    txt(sl,t['title_main'],0.7,1.0,8.6,0.75,size=36,bold=True,color=WHITE)
    txt(sl,t['title_sub'], 0.7,1.82,8.6,0.5, size=18,color=ICE)
    txt(sl,t['title_geo'], 0.7,2.38,8.6,0.38,size=12,color=LGRAY)
    cfr_f=cfr_v(D['fd_tot'],D['ftot']); cfr_m=cfr_v(D['md_tot'],D['mtot'])
    cfr_t=cfr_v(D['total_d'],D['total'])
    tiles=[(t['total_cases'],f"{D['total']:,}",""),
           (t['total_deaths'],f"{D['total_d']:,}",f"CFR {cfr_t}%"),
           (t['female_cases'],f"{D['ftot']:,}",f"CFR {cfr_f}%"),
           (t['male_cases'],  f"{D['mtot']:,}",f"CFR {cfr_m}%")]
    for i,(lbl,val,sub) in enumerate(tiles):
        x=0.3+i*2.33
        rect(sl,x,2.9,2.22,1.08,rgb("092540"),radius=True,line=rgb("1E5080"))
        txt(sl,lbl,x+0.1,2.97,2.02,0.24,size=8.5,color=LGRAY,align=PP_ALIGN.CENTER)
        txt(sl,val,x+0.1,3.22,2.02,0.44,size=24,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        if sub: txt(sl,sub,x+0.1,3.68,2.02,0.22,size=9,color=rgb("5EADD4"),align=PP_ALIGN.CENTER)
    gtiles=[(str(D['np']),t['provinces']),(str(D['nz']),t['health_zones']),(str(D['na']),t['health_areas'])]
    for i,(val,lbl) in enumerate(gtiles):
        x=1.15+i*2.9
        rect(sl,x,4.14,2.22,0.88,rgb("0A1628"),radius=True,line=rgb("1E5080"))
        txt(sl,val,x+0.08,4.19,2.06,0.38,size=22,bold=True,color=rgb("7DCFDF"),align=PP_ALIGN.CENTER)
        txt(sl,lbl,x+0.08,4.58,2.06,0.28,size=9,color=LGRAY,align=PP_ALIGN.CENTER)
    txt(sl,f"{t['data_date']}: {D['ddate']}  |  {t['earliest_onset']}: {D['earliest']}  |  {t['source']}",
        0.5,5.26,9.0,0.22,size=8.5,color=LGRAY)

def s2_cases(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s2_title'],f"{t['data_date']}: {D['ddate']}")
    labs=D['labels']+[t['unknown']]
    fc=D['fc']+[D['fc_unk']]; mc=D['mc']+[D['mc_unk']]; uc=D['uc']+[D['uc_unk']]
    series=[(t['female'],labs,fc),(t['male'],labs,mc)]
    if sum(uc)>0: series.append((t['unknown'],labs,uc))
    chart(sl,XL_CHART_TYPE.COLUMN_CLUSTERED,series,0.4,1.1,9.2,4.0,colors=[FEM_C,MAL_C,UNK_SEX])
    foot(sl,f"Unknown age: {D['dq']['unk_age']} cases included in totals.")

def s3_deaths(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s3_title'],f"{t['data_date']}: {D['ddate']}")
    labs=D['labels']+[t['unknown']]
    fd=D['fd']+[D['fd_unk']]; md=D['md']+[D['md_unk']]; ud=D['ud']+[D['ud_unk']]
    series=[(t['female_d'],labs,fd),(t['male_d'],labs,md)]
    if sum(ud)>0: series.append((t['unknown_d'],labs,ud))
    chart(sl,XL_CHART_TYPE.COLUMN_CLUSTERED,series,0.4,1.1,9.2,3.6,colors=[rgb("5B8DB8"),rgb("D4875A"),UNK_SEX])
    txt(sl,t['death_note'],0.5,4.75,9.0,0.22,size=8,color=LGRAY,italic=True)
    foot(sl,f"Unknown age: {D['dq']['unk_age_d']} deaths included in totals.")

def s4_cfr(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s4_title'],f"{t['data_date']}: {D['ddate']}")

    labels   = D['labels']
    cfr_vals = D['cfr_age']

    # Split into three series by severity tier so each renders in its own colour.
    # python-pptx cannot colour individual bars within a single series;
    # separate series with None placeholders for non-applicable bars is the fix.
    crit = [v if v>=40  else None for v in cfr_vals]  # red
    high = [v if 25<=v<40 else None for v in cfr_vals]  # amber
    mod  = [v if v<25  else None for v in cfr_vals]  # teal

    # Replace None with 0 for the chart data (zero-height bar = invisible),
    # but keep track of which bars actually have values for label placement.
    def nz(lst): return [v if v is not None else 0 for v in lst]

    from pptx.chart.data import CategoryChartData
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series(t['cfr_crit'], tuple(nz(crit)))
    cd.add_series(t['cfr_high'], tuple(nz(high)))
    cd.add_series(t['cfr_mod'],  tuple(nz(mod)))
    cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                             I(0.4),I(1.1),I(9.2),I(3.6), cd)
    ch = cf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.TOP
    ch.legend.include_in_layout = False
    # Colours and data labels — show value only for non-zero bars
    for s_idx, (series, sentinel_list, col) in enumerate(
            zip(ch.series, [crit, high, mod], [CORAL, AMBER, TEAL])):
        f = series.format.fill; f.solid(); f.fore_color.rgb = col
        series.data_labels.show_value = True
        series.data_labels.font.size = P(8)
        # Zero-out labels on placeholder bars by setting number format
        # to suppress zeros: use a custom number format "0.#;;;" 
        # (positive;negative;zero;text — zero format is empty)
        from pptx.oxml.ns import qn
        from lxml import etree
        dLbls = series._element.find(qn('c:dLbls'))
        if dLbls is not None:
            numFmt = etree.SubElement(dLbls, qn('c:numFmt'))
            numFmt.set('formatCode', '0.#;;;')
            numFmt.set('sourceLinked', '0')
    # Y-axis title
    va = ch.value_axis; va.has_title = True
    va.axis_title.text_frame.text = t['cfr_pct']
    va.axis_title.text_frame.paragraphs[0].runs[0].font.size = P(10)
    ch.value_axis.maximum_scale = 75
    ch.value_axis.minimum_scale = 0

    # The legend from the chart itself now shows correct colours — no manual legend needed.
    unk=D['dq']['unk_age']
    foot(sl,f"CFR shown for {D['total']-unk} cases with known age. {unk} unknown-age cases excluded from this chart.")


def s5_div(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,NAVY_D)
    txt(sl,t['s5_title'],0.7,1.5,8.6,0.8,size=36,bold=True,color=WHITE)
    txt(sl,t['s5_sub'],  0.7,2.4,8.6,0.85,size=16,color=ICE)
    txt(sl,f"{t['data_date']}: {D['ddate']}  |  {D['total']:,} {t['cases'].lower()}  |  "
           f"{D['total_d']} {t['deaths'].lower()}  |  "
           f"{t['earliest_onset']}: {D['earliest']}  |  {t['source']}",
        0.7,5.22,8.6,0.25,size=8.5,color=LGRAY)

def s6_global(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    gc,gd=D['gc'],D['gd']
    hdr(sl,t['s6_title'],
        f"{t['data_date']}: {D['ddate']}  |  {D['total']:,} {t['cases'].lower()}  |  {D['total_d']} {t['deaths'].lower()}")
    # Cases chart: Child/Adult/Unknown × Female/Male
    genders=[t['female'],t['male']]
    series_c=[(t['child'],   genders,[gc['cf'],gc['cm']]),
               (t['adult'],  genders,[gc['af'],gc['am']]),
               (t['unk_age'],genders,[gc['uf'],gc['um']])]
    chart(sl,XL_CHART_TYPE.COLUMN_STACKED,series_c,0.3,1.1,4.5,3.8,
          colors=[CHILD_C,ADULT_C,UNK_C])
    series_d=[(t['child_d'],  genders,[gd['cf'],gd['cm']]),
               (t['adult_d'], genders,[gd['af'],gd['am']]),
               (t['unk_age_d'],genders,[gd['uf'],gd['um']])]
    chart(sl,XL_CHART_TYPE.COLUMN_STACKED,series_d,5.2,1.1,4.5,3.8,
          colors=[CHILD_D,ADULT_D,UNK_C])
    cfr_c=cfr_v(gd['ct'],gc['ct']); cfr_a=cfr_v(gd['at'],gc['at'])
    foot(sl,f"{t['child']}: {gc['ct']} {t['cases'].lower()} (CFR {cfr_c}%)  |  "
            f"{t['adult']}: {gc['at']} (CFR {cfr_a}%)  |  "
            f"Unknown age: {gc['unk']} cases included in totals")

def s7_metrics(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s7_title'],f"{t['data_date']}: {D['ddate']}")
    gc,gd=D['gc'],D['gd']
    cfr_c=cfr_v(gd['ct'],gc['ct']); cfr_a=cfr_v(gd['at'],gc['at'])
    cards=[
        (0.3, 1.1, f"{t['child']} — {t['cases'].lower()}",f"{gc['ct']}",
         f"{pct_v(gc['ct'],gc['total'])}% of total",rgb("FDECEA")),
        (2.65,1.1, f"{t['female']} / {t['male']}",f"{gc['cf']} / {gc['cm']}",
         f"{pct_v(gc['cf'],gc['ct'])}% F",rgb("FDECEA")),
        (5.0, 1.1, f"{t['adult']} — {t['cases'].lower()}",f"{gc['at']}",
         f"{pct_v(gc['at'],gc['total'])}% of total",rgb("E0F2F4")),
        (7.35,1.1, f"{t['female']} / {t['male']}",f"{gc['af']} / {gc['am']}",
         f"{pct_v(gc['af'],gc['at'])}% F",rgb("E0F2F4")),
        (0.3, 2.55,f"{t['child']} — {t['deaths'].lower()}",f"{gd['ct']}",f"CFR {cfr_c}%",rgb("FDECEA")),
        (2.65,2.55,f"{t['female']} / {t['male']}",f"{gd['cf']} / {gd['cm']}","",rgb("FDECEA")),
        (5.0, 2.55,f"{t['adult']} — {t['deaths'].lower()}",f"{gd['at']}",f"CFR {cfr_a}%",rgb("E0F2F4")),
        (7.35,2.55,f"{t['female']} / {t['male']}",f"{gd['af']} / {gd['am']}","",rgb("E0F2F4")),
    ]
    for x,y,lbl,val,sub,bg in cards:
        rect(sl,x,y,2.2,1.25,bg,line=rgb("DDE3EA"))
        txt(sl,lbl,x+0.1,y+0.08,2.0,0.25,size=8.5,color=LGRAY,align=PP_ALIGN.CENTER)
        txt(sl,val,x+0.1,y+0.35,2.0,0.45,size=20,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        if sub: txt(sl,sub,x+0.1,y+0.85,2.0,0.28,size=9,color=TEAL,align=PP_ALIGN.CENTER)
    txt(sl,t['cfr_comparison'],0.5,4.05,3.5,0.28,size=11,bold=True,color=DARK)
    cfr_ov=cfr_v(D['total_d'],D['total'])
    for i,(lbl,val,col) in enumerate([(t['child'],cfr_c,CHILD_C),(t['adult'],cfr_a,ADULT_C),(t['overall'],cfr_ov,LGRAY)]):
        bw=(val/60)*6.5; yy=4.42+i*0.32
        txt(sl,lbl,0.5,yy,1.8,0.26,size=10,color=DARK)
        rect(sl,2.4,yy+0.03,bw,0.2,col)
        txt(sl,f"{val}%",2.4+bw+0.08,yy,0.8,0.25,size=10,bold=True,color=col)
    foot(sl,f"{t['child']} {t['cfr']} {cfr_c}% vs {t['adult']} {cfr_a}%  |  Unknown age: {gc['unk']} cases included in totals.")


def s89_prov(prs,D,t,is_d):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s9_title'] if is_d else t['s8_title'],
        t['confirmed_deaths'] if is_d else t['confirmed_only'])
    pnames=[p['name'] for p in D['pdata']]
    src=[p['deaths'] if is_d else p['cases'] for p in D['pdata']]
    # 100% stacked: Child F, Child M, Adult F, Adult M, Unknown F, Unknown M
    series=[
        (t['child_f'],pnames,[s['cf'] for s in src]),
        (t['child_m'],pnames,[s['cm'] for s in src]),
        (t['adult_f'],pnames,[s['af'] for s in src]),
        (t['adult_m'],pnames,[s['am'] for s in src]),
        (t['unk_f'],  pnames,[s['uf'] for s in src]),
        (t['unk_m'],  pnames,[s['um'] for s in src]),
    ]
    chart(sl,XL_CHART_TYPE.COLUMN_STACKED_100,series,0.3,1.05,9.4,2.72,
          colors=[rgb("E8A89C"),CHILD_C,rgb("7BBFC8"),ADULT_C,rgb("D1D5DB"),UNK_C])
    HB=NAVY_D; HW=WHITE
    CW=[1.42,0.78,0.78,0.82,0.78,0.78,0.82,0.72,0.62,0.67]
    hrow=[(t['province'],True,HW,HB,7.5,PP_ALIGN.LEFT),
          (t['child_f'], True,HW,HB,7.5,PP_ALIGN.CENTER),(t['child_m'],True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['child_t'], True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['adult_f'], True,HW,HB,7.5,PP_ALIGN.CENTER),(t['adult_m'],True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['adult_t'], True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['unk_t'],   True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['total'],   True,HW,HB,7.5,PP_ALIGN.CENTER),
          (t['cfr'],     True,HW,HB,7.5,PP_ALIGN.CENTER)]
    rows=[hrow]
    for i,(p,s) in enumerate(zip(D['pdata'],src)):
        bg=rgb("F0F4F8") if i%2==0 else WHITE
        cv=cfr_v(p['deaths']['total'],p['cases']['total']); cc=cfr_col(cv)
        def cell(v,tot,bold=False,col=None):
            return (f"{pct_v(v,tot)}%\n({v})",bold,col or DARK,bg,7,PP_ALIGN.CENTER)
        rows.append([
            (p['name'],True,DARK,bg,8,PP_ALIGN.LEFT),
            cell(s['cf'],s['total']),cell(s['cm'],s['total']),cell(s['ct'],s['total'],True,CHILD_C),
            cell(s['af'],s['total']),cell(s['am'],s['total']),cell(s['at'],s['total'],True,ADULT_C),
            cell(s['unk'],s['total'],False,UNK_C),
            (str(s['total']),True,DARK,bg,8,PP_ALIGN.CENTER),(f"{cv}%",True,cc,bg,8,PP_ALIGN.CENTER),
        ])
    tbl(sl,rows,0.3,3.93,9.4,1.12,cw=CW)
    total_unk=sum(s['unk'] for s in src)
    note=t['prov_note']
    if total_unk>0: note+=f"  Unknown age: {total_unk} records shown."
    txt(sl,note,0.3,5.12,9.4,0.2,size=7.5,color=LGRAY,italic=True)

def s10_zone_cases(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s10_title'],t['s10_sub'])
    zn=[z['name'] for z in D['zdata']]
    series=[(t['child'],  zn,[z['cases']['ct']  for z in D['zdata']]),
             (t['adult'], zn,[z['cases']['at']  for z in D['zdata']]),
             (t['unk_age'],zn,[z['cases']['unk'] for z in D['zdata']])]
    chart(sl,XL_CHART_TYPE.COLUMN_STACKED,series,0.3,1.05,9.4,3.9,colors=[CHILD_C,ADULT_C,UNK_C])
    foot(sl,D['zone_gender_note'])

def s11_zone_deaths(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s11_title'],t['death_note'])
    zn=[z['name'] for z in D['zdata']]
    series=[(t['child_d'], zn,[z['deaths']['ct']  for z in D['zdata']]),
             (t['adult_d'],zn,[z['deaths']['at']  for z in D['zdata']]),
             (t['unk_age_d'],zn,[z['deaths']['unk'] for z in D['zdata']])]
    chart(sl,XL_CHART_TYPE.COLUMN_STACKED,series,0.3,1.05,5.8,3.5,colors=[CHILD_D,ADULT_D,UNK_C])
    HB=NAVY_D
    rows=[[(t['zone'],True,WHITE,HB,8.5,PP_ALIGN.LEFT),(t['cases'],True,WHITE,HB,8.5,PP_ALIGN.CENTER),
           (t['deaths'],True,WHITE,HB,8.5,PP_ALIGN.CENTER),(t['cfr'],True,WHITE,HB,8.5,PP_ALIGN.CENTER)]]
    for i,z in enumerate(D['zdata']):
        bg=rgb("F4F6F8") if i%2==0 else WHITE
        cv=cfr_v(z['deaths']['total'],z['cases']['total']); cc=cfr_col(cv)
        rows.append([(z['name'],False,DARK,bg,8.5,PP_ALIGN.LEFT),
                     (str(z['cases']['total']),False,DARK,bg,8.5,PP_ALIGN.CENTER),
                     (str(z['deaths']['total']),False,DARK,bg,8.5,PP_ALIGN.CENTER),
                     (f"{cv}%",True,cc,bg,8.5,PP_ALIGN.CENTER)])
    tbl(sl,rows,6.4,1.15,3.3,3.4,cw=[1.35,0.7,0.75,0.5])
    foot(sl,D['zone_deaths_note'])

def s12_zone_gender(prs,D,t):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl,SLATE)
    hdr(sl,t['s12_title'],f"{t['data_date']}: {D['ddate']}")
    zn=[z['name'] for z in D['zdata']]
    fc_z=[z['cases']['cf']+z['cases']['af'] for z in D['zdata']]
    mc_z=[z['cases']['cm']+z['cases']['am'] for z in D['zdata']]
    fd_z=[z['deaths']['cf']+z['deaths']['af'] for z in D['zdata']]
    md_z=[z['deaths']['cm']+z['deaths']['am'] for z in D['zdata']]
    series=[(t['female_cases_s12'],zn,fc_z),(t['male_cases_s12'],zn,mc_z),
             (t['female_deaths_s12'],zn,fd_z),(t['male_deaths_s12'],zn,md_z)]
    chart(sl,XL_CHART_TYPE.COLUMN_CLUSTERED,series,0.3,1.05,9.4,4.0,
          colors=[FEM_C,MAL_C,rgb("A0B8D8"),rgb("EAB48A")])
    foot(sl,D['zone_gender_note'])

# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    ap=argparse.ArgumentParser(description="BVD Sex and Age Analysis MVE report v4.0")
    ap.add_argument("input"); ap.add_argument("--lang",choices=["en","fr"],default="en")
    ap.add_argument("--out",default=None)
    args=ap.parse_args()
    fp=Path(args.input)
    if not fp.exists(): print(f"ERROR: {fp} not found",file=sys.stderr); sys.exit(1)
    t=T[args.lang]
    print(f"\n{'='*62}\nBVD Sex and Age Analysis v4.0  |  Lang: {args.lang.upper()}\nInput: {fp.name}\n{'='*62}")
    print("\n[1/3] Extracting data...")
    D,conf,cd=extract(str(fp))
    print(f"  Confirmed: {D['total']:,}  |  Deaths: {D['total_d']}  |  CFR {cfr_v(D['total_d'],D['total'])}%")
    print(f"  Unknown age: {D['dq']['unk_age']} cases ({D['dq']['unk_age_pct']}%), {D['dq']['unk_age_d']} deaths")
    print(f"  Unknown sex: {D['dq']['unk_sex']} cases")
    print(f"  Invalid age codes found: {D['dq']['inv_codes']}")
    print("\n[2/3] Running quality gates...")
    errors,warnings=run_gates(conf,cd,D['gc'],D['gd'],D['pdata'],D)
    D['dq']['reconciled']=(len(errors)==0)
    for w in warnings: print(f"  ⚠️  {w}")
    for e in errors:   print(f"  ❌ {e}")
    if not errors and not warnings: print("  ✅ All quality gates passed.")
    elif not errors: print(f"  ✅ Passed with {len(warnings)} warning(s).")
    else:
        print(f"\n❌ {len(errors)} gate(s) failed — fix before generating slides."); sys.exit(2)
    print("\n[3/3] Generating 13 slides (incl. DQ slide)...")
    prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
    s1_title(prs,D,t)        #  1 – Title
    s0_dq(prs,D,t)           #  2 – Data quality
    s2_cases(prs,D,t)        #  3 – Cases by age/gender
    s3_deaths(prs,D,t)       #  4 – Deaths by age/gender
    s4_cfr(prs,D,t)          #  5 – CFR
    s5_div(prs,D,t)          #  6 – Section divider
    s6_global(prs,D,t)       #  7 – Global child/adult
    s7_metrics(prs,D,t)      #  8 – Key metrics
    s89_prov(prs,D,t,False)  #  9 – Province cases
    s89_prov(prs,D,t,True)   # 10 – Province deaths
    s10_zone_cases(prs,D,t)  # 11 – Zone cases
    s11_zone_deaths(prs,D,t) # 12 – Zone deaths + CFR
    s12_zone_gender(prs,D,t) # 13 – Zone gender
    out=args.out or f"BVD-Sex-Age-Analysis-Report-{D['dstr']}.pptx"
    prs.save(out)
    print(f"\n✅ Saved: {out}  ({13} slides)")
    if warnings: print(f"   ⚠️  {len(warnings)} warning(s) — review footnotes on affected slides.")

if __name__=="__main__":
    main()